import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches

df = pd.read_csv('data/mpg.csv')
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
title_only_slide_layout = prs.slide_layouts[5]

slide1 = prs.slides.add_slide(title_slide_layout)

title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "Trying out PowerPoint Automation"
subtitle.text = "With python-pptx and GitHub action!"

slide2 = prs.slides.add_slide(title_only_slide_layout)
slide2.shapes.title.text = 'Add Image with Python'

img = sns.heatmap(df.corr(), annot = True).get_figure()
img.savefig('graph/heatmap1.png')

left = Inches(3)
top = Inches(4)
height = Inches(4)
pic = slide2.shapes.add_picture('heatmap1.png', left, top, height=height)

prs.save('report.pptx')